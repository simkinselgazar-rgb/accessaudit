"""Office document capture for WCAG accessibility testing.

Handles DOCX, XLSX, and PPTX files, extracting content into
CaptureData with pseudo-HTML for downstream WCAG checks.
"""
from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

from models import CaptureData

logger = logging.getLogger(__name__)


# ═════════════════════════════════════════════════════════════════════════════
#  DOCX capture
# ═════════════════════════════════════════════════════════════════════════════

def capture_docx(file_path: str, review_dir: str) -> CaptureData:
    """Capture a DOCX document for WCAG accessibility testing.

    Extracts headings, links, images (with alt text), and tables,
    then builds pseudo-HTML for check compatibility.

    Args:
        file_path: Path to the DOCX file.
        review_dir: Path to the review output directory.

    Returns:
        A populated CaptureData instance.
    """
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    captures_dir = os.path.join(review_dir, "captures")
    os.makedirs(captures_dir, exist_ok=True)

    capture_data = CaptureData(
        file_path=file_path,
        file_type="docx",
        review_dir=review_dir,
        captures_dir=captures_dir,
    )

    try:
        doc = Document(file_path)
    except Exception:
        logger.exception("Failed to open DOCX: %s", file_path)
        return capture_data

    try:
        # Title from core properties or filename
        title = ""
        try:
            title = doc.core_properties.title or ""
        except Exception:
            pass  # default to filename-derived title if core_properties is missing/corrupt
        if not title:
            title = Path(file_path).stem
        capture_data.title = title

        # ── Document language ───────────────────────────────────────
        doc_language = "en"
        try:
            from docx.oxml.ns import qn as _qn
            # Try default paragraph style language
            try:
                from docx.enum.style import WD_STYLE_TYPE
                default_style = doc.styles.default(WD_STYLE_TYPE.PARAGRAPH)
            except Exception:
                default_style = None
            if default_style is not None and hasattr(default_style, "font") and default_style.font.language:
                lang_obj = default_style.font.language
                if lang_obj.latin:
                    doc_language = str(lang_obj.latin)
            if doc_language == "en":
                # Try settings.xml themeFontLang
                settings_elem = getattr(doc, "settings", None)
                if settings_elem is not None:
                    s_elem = getattr(settings_elem, "element", None)
                    if s_elem is not None:
                        theme_font_lang = s_elem.find(_qn("w:themeFontLang"))
                        if theme_font_lang is not None:
                            val = theme_font_lang.get(_qn("w:val"), "")
                            if val:
                                doc_language = val
            if doc_language == "en":
                # Try first few paragraph runs
                for para in doc.paragraphs:
                    for run in para.runs:
                        if hasattr(run.font, "language") and run.font.language:
                            if run.font.language.latin:
                                doc_language = str(run.font.language.latin)
                                break
                    if doc_language != "en":
                        break
            if doc_language == "en":
                # Try w:lang elements in styles.xml
                styles_elem = doc.styles.element
                lang_elems = styles_elem.findall(".//" + _qn("w:lang"))
                for lang_elem in lang_elems:
                    val = lang_elem.get(_qn("w:val"), "")
                    if val:
                        doc_language = val
                        break
        except Exception:
            logger.debug("Failed to extract document language from DOCX")
        capture_data.user_context["doc_language"] = doc_language

        headings: list[dict] = []
        links: list[dict] = []
        images: list[dict] = []
        tables_data: list[dict] = []
        body_html_parts: list[str] = []

        # ── Paragraphs ──────────────────────────────────────────────
        for para in doc.paragraphs:
            style_name = (para.style.name or "").lower() if para.style else ""
            text = para.text.strip()

            # Headings
            if style_name.startswith("heading"):
                try:
                    level = int(style_name.replace("heading", "").strip())
                except (ValueError, IndexError):
                    level = 1
                level = max(1, min(level, 6))
                headings.append({
                    "tag": f"h{level}",
                    "level": level,
                    "text": text,
                })
                body_html_parts.append(f"<h{level}>{_escape(text)}</h{level}>")
            elif text:
                body_html_parts.append(f"<p>{_escape(text)}</p>")

            # Links in runs
            for run in para.runs:
                try:
                    for rel in run.part.rels.values():
                        if "hyperlink" in str(rel.reltype).lower():
                            links.append({
                                "text": run.text.strip(),
                                "href": rel.target_ref or "",
                            })
                except Exception:
                    pass  # best-effort — skip run if its relationship table is unreadable

        # ── Hyperlinks from relationships ────────────────────────────
        try:
            for rel_id, rel in doc.part.rels.items():
                if "hyperlink" in str(rel.reltype).lower():
                    if rel.target_ref and rel.target_ref not in [l.get("href") for l in links]:
                        links.append({
                            "text": "",
                            "href": rel.target_ref,
                        })
        except Exception:
            logger.debug("Relationship-based hyperlink extraction failed")

        # ── Images ───────────────────────────────────────────────────
        img_dir = os.path.join(captures_dir, "images")
        os.makedirs(img_dir, exist_ok=True)
        img_index = 0

        for rel_id, rel in doc.part.rels.items():
            try:
                if "image" in str(rel.reltype).lower():
                    image_part = rel.target_part
                    ext = Path(image_part.partname).suffix or ".png"
                    img_filename = f"docx_img_{img_index}{ext}"
                    img_path = os.path.join(img_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image_part.blob)

                    # Try to find alt text from inline shapes
                    alt_text = _find_docx_image_alt(doc, rel_id)

                    images.append({
                        "path": img_path,
                        "alt": alt_text,
                        "width": 0,
                        "height": 0,
                    })
                    alt_attr = _escape(alt_text)
                    body_html_parts.append(
                        f'<img src="{_escape(img_path)}" alt="{alt_attr}">'
                    )
                    img_index += 1
            except Exception:
                logger.debug("Failed to extract image from relationship %s", rel_id)

        # ── Tables ───────────────────────────────────────────────────
        for tbl_idx, table in enumerate(doc.tables):
            try:
                rows_data: list[list[str]] = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_cells)

                tables_data.append({
                    "index": tbl_idx,
                    "rows": rows_data,
                    "rowCount": len(rows_data),
                    "columnCount": len(rows_data[0]) if rows_data else 0,
                    "caption": "",
                })

                # Build HTML table
                body_html_parts.append("<table>")
                for r_idx, row_cells in enumerate(rows_data):
                    body_html_parts.append("  <tr>")
                    cell_tag = "th" if r_idx == 0 else "td"
                    for cell in row_cells:
                        body_html_parts.append(
                            f"    <{cell_tag}>{_escape(cell)}</{cell_tag}>"
                        )
                    body_html_parts.append("  </tr>")
                body_html_parts.append("</table>")
            except Exception:
                logger.debug("Failed to extract table %d", tbl_idx)

        # ── Build pseudo-HTML ────────────────────────────────────────
        pseudo_html = _wrap_html(title, body_html_parts, links, lang=doc_language)

        capture_data.html = pseudo_html
        capture_data.headings = headings
        capture_data.links = links
        capture_data.images = images
        capture_data.tables = tables_data

    except Exception:
        logger.exception("DOCX capture failed for %s", file_path)

    return capture_data


def _find_docx_image_alt(doc: Any, rel_id: str) -> str:
    """Attempt to extract alt text for an image identified by rel_id."""
    try:
        from docx.oxml.ns import qn
        for para in doc.paragraphs:
            for run in para.runs:
                drawings = run._element.findall(qn("w:drawing"))
                for drawing in drawings:
                    # Check inline images
                    for inline in drawing.findall(qn("wp:inline")):
                        doc_pr = inline.find(qn("wp:docPr"))
                        if doc_pr is not None:
                            descr = doc_pr.get("descr", "")
                            name = doc_pr.get("name", "")
                            blip_fill = inline.find(".//" + qn("a:blip"))
                            if blip_fill is not None:
                                embed = blip_fill.get(qn("r:embed"), "")
                                if embed == rel_id:
                                    return descr or name
                    # Check anchored images
                    for anchor in drawing.findall(qn("wp:anchor")):
                        doc_pr = anchor.find(qn("wp:docPr"))
                        if doc_pr is not None:
                            descr = doc_pr.get("descr", "")
                            name = doc_pr.get("name", "")
                            blip_fill = anchor.find(".//" + qn("a:blip"))
                            if blip_fill is not None:
                                embed = blip_fill.get(qn("r:embed"), "")
                                if embed == rel_id:
                                    return descr or name
    except Exception:
        pass  # default to empty alt text if anchor XML cannot be parsed
    return ""


# ═════════════════════════════════════════════════════════════════════════════
#  XLSX capture
# ═════════════════════════════════════════════════════════════════════════════

def capture_xlsx(file_path: str, review_dir: str) -> CaptureData:
    """Capture an XLSX spreadsheet for WCAG accessibility testing.

    Extracts each sheet as a table and builds pseudo-HTML.

    Args:
        file_path: Path to the XLSX file.
        review_dir: Path to the review output directory.

    Returns:
        A populated CaptureData instance.
    """
    from openpyxl import load_workbook

    captures_dir = os.path.join(review_dir, "captures")
    os.makedirs(captures_dir, exist_ok=True)

    capture_data = CaptureData(
        file_path=file_path,
        file_type="xlsx",
        review_dir=review_dir,
        captures_dir=captures_dir,
    )

    try:
        wb = load_workbook(file_path, data_only=True, read_only=False)
    except Exception:
        logger.exception("Failed to open XLSX: %s", file_path)
        return capture_data

    try:
        capture_data.title = wb.properties.title or Path(file_path).stem

        tables_data: list[dict] = []
        body_html_parts: list[str] = []
        cell_a1_content: dict[str, Any] = {}

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            try:
                ws = wb[sheet_name]
                rows_data: list[list[str]] = []

                for row in ws.iter_rows(values_only=True):
                    row_cells = [str(cell) if cell is not None else "" for cell in row]
                    rows_data.append(row_cells)

                # Extract merged cell ranges
                merged_ranges: list[str] = []
                try:
                    for merge_range in ws.merged_cells.ranges:
                        merged_ranges.append(str(merge_range))
                except Exception:
                    logger.debug("Failed to extract merged cells for sheet '%s'", sheet_name)

                # Extract cell A1 content
                try:
                    a1_val = ws["A1"].value
                    cell_a1_content[sheet_name] = str(a1_val) if a1_val is not None else ""
                except Exception:
                    logger.debug("Failed to read A1 for sheet '%s'", sheet_name)

                tables_data.append({
                    "index": sheet_idx,
                    "caption": sheet_name,
                    "rows": rows_data,
                    "rowCount": len(rows_data),
                    "columnCount": len(rows_data[0]) if rows_data else 0,
                    "merged_cells": merged_ranges,
                })

                # Build HTML table for this sheet
                body_html_parts.append(f"<h2>{_escape(sheet_name)}</h2>")
                body_html_parts.append(f"<table><caption>{_escape(sheet_name)}</caption>")
                for r_idx, row_cells in enumerate(rows_data):
                    body_html_parts.append("  <tr>")
                    cell_tag = "th" if r_idx == 0 else "td"
                    for cell in row_cells:
                        body_html_parts.append(
                            f"    <{cell_tag}>{_escape(cell)}</{cell_tag}>"
                        )
                    body_html_parts.append("  </tr>")
                body_html_parts.append("</table>")
            except Exception:
                logger.debug("Failed to process sheet '%s'", sheet_name)

        # Build headings from sheet names
        headings = [
            {"tag": "h2", "level": 2, "text": name}
            for name in wb.sheetnames
        ]

        capture_data.user_context["cell_a1_content"] = cell_a1_content

        pseudo_html = _wrap_html(capture_data.title, body_html_parts)
        capture_data.html = pseudo_html
        capture_data.headings = headings
        capture_data.tables = tables_data

    except Exception:
        logger.exception("XLSX capture failed for %s", file_path)
    finally:
        wb.close()

    return capture_data


# ═════════════════════════════════════════════════════════════════════════════
#  PPTX capture
# ═════════════════════════════════════════════════════════════════════════════

def capture_pptx(file_path: str, review_dir: str) -> CaptureData:
    """Capture a PPTX presentation for WCAG accessibility testing.

    Extracts slides, text, images, and tables, then builds pseudo-HTML.

    Args:
        file_path: Path to the PPTX file.
        review_dir: Path to the review output directory.

    Returns:
        A populated CaptureData instance.
    """
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    captures_dir = os.path.join(review_dir, "captures")
    os.makedirs(captures_dir, exist_ok=True)

    capture_data = CaptureData(
        file_path=file_path,
        file_type="pptx",
        review_dir=review_dir,
        captures_dir=captures_dir,
    )

    try:
        prs = Presentation(file_path)
    except Exception:
        logger.exception("Failed to open PPTX: %s", file_path)
        return capture_data

    try:
        # Title from core properties or filename
        title = ""
        try:
            title = prs.core_properties.title or ""
        except Exception:
            pass  # default to filename-derived title if core_properties is missing/corrupt
        if not title:
            title = Path(file_path).stem
        capture_data.title = title

        # ── Presentation language ───────────────────────────────────
        doc_language = "en"
        try:
            from lxml import etree
            # Try presentation.xml defaultTextStyle > defRPr lang
            prs_elem = prs.presentation._element if hasattr(prs, "presentation") else prs._element
            nsmap = {
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }
            def_text_style = prs_elem.find(".//p:defaultTextStyle", nsmap)
            if def_text_style is not None:
                for child in def_text_style:
                    def_rpr = child.find("a:defRPr", nsmap)
                    if def_rpr is not None:
                        lang = def_rpr.get("lang", "")
                        if lang:
                            doc_language = lang
                            break
            if doc_language == "en":
                # Try first slide text runs for lang attribute
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    try:
                                        rpr = run._r.find(
                                            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                                        )
                                        if rpr is not None:
                                            lang = rpr.get("lang", "")
                                            if lang:
                                                doc_language = lang
                                                raise StopIteration
                                    except StopIteration:
                                        raise
                                    except Exception:
                                        pass  # best-effort — skip rPr if its lang attribute is unreadable
                    # Continue checking all slides until language found
        except StopIteration:
            pass
        except Exception:
            logger.debug("Failed to extract presentation language from PPTX")
        capture_data.user_context["doc_language"] = doc_language

        # ── Slide transitions (auto-advance timing) ─────────────────
        slide_transitions: list[dict] = []
        try:
            nsmap_p = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }
            for s_idx, slide in enumerate(prs.slides):
                try:
                    slide_elem = slide._element
                    transition = slide_elem.find("p:transition", nsmap_p)
                    if transition is not None:
                        adv_tm = transition.get("advTm")
                        adv_click = transition.get("advClick", "1")
                        auto_advance_ms = int(adv_tm) if adv_tm else None
                        advance_on_click = adv_click != "0"
                        slide_transitions.append({
                            "slide": s_idx + 1,
                            "auto_advance_ms": auto_advance_ms,
                            "advance_on_click": advance_on_click,
                        })
                    else:
                        slide_transitions.append({
                            "slide": s_idx + 1,
                            "auto_advance_ms": None,
                            "advance_on_click": True,
                        })
                except Exception:
                    logger.debug("Failed to extract transition for slide %d", s_idx + 1)
        except Exception:
            logger.debug("Failed to extract slide transitions from PPTX")
        capture_data.user_context["slide_transitions"] = slide_transitions

        headings: list[dict] = []
        images: list[dict] = []
        tables_data: list[dict] = []
        body_html_parts: list[str] = []
        links: list[dict] = []

        img_dir = os.path.join(captures_dir, "images")
        os.makedirs(img_dir, exist_ok=True)
        img_index = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            body_html_parts.append(
                f'<section aria-label="Slide {slide_num}">'
            )

            slide_title = ""
            for shape in slide.shapes:
                try:
                    # ── Title shapes ─────────────────────────────────
                    if shape.has_text_frame:
                        if shape.shape_type is not None and hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format is not None:
                                ph_idx = shape.placeholder_format.idx
                                # idx 0 = title, idx 1 = subtitle
                                if ph_idx == 0:
                                    slide_title = shape.text.strip()

                    # ── Text frames ──────────────────────────────────
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            # Check for hyperlinks
                            for run in para.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    links.append({
                                        "text": run.text.strip(),
                                        "href": run.hyperlink.address,
                                        "slide": slide_num,
                                    })
                            body_html_parts.append(f"  <p>{_escape(text)}</p>")

                    # ── Images / pictures ────────────────────────────
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            ext = image.content_type.split("/")[-1] if image.content_type else "png"
                            if ext == "jpeg":
                                ext = "jpg"
                            img_filename = f"pptx_slide{slide_num}_img{img_index}.{ext}"
                            img_path = os.path.join(img_dir, img_filename)
                            with open(img_path, "wb") as f:
                                f.write(image.blob)

                            alt_text = shape.name or ""
                            try:
                                if hasattr(shape, "_element"):
                                    from pptx.oxml.ns import qn
                                    nvPicPr = shape._element.find(qn("p:nvPicPr"))
                                    if nvPicPr is not None:
                                        cNvPr = nvPicPr.find(qn("p:cNvPr"))
                                        if cNvPr is not None:
                                            alt_text = cNvPr.get("descr", "") or cNvPr.get("name", "")
                            except Exception:
                                pass  # default to existing alt_text if cNvPr XML cannot be parsed

                            images.append({
                                "path": img_path,
                                "alt": alt_text,
                                "slide": slide_num,
                                "width": shape.width,
                                "height": shape.height,
                            })
                            body_html_parts.append(
                                f'  <img src="{_escape(img_path)}" alt="{_escape(alt_text)}">'
                            )
                            img_index += 1
                        except Exception:
                            logger.debug("Failed to extract image from slide %d", slide_num)

                    # ── Tables ───────────────────────────────────────
                    if shape.has_table:
                        try:
                            table = shape.table
                            rows_data: list[list[str]] = []
                            for row in table.rows:
                                row_cells = [cell.text.strip() for cell in row.cells]
                                rows_data.append(row_cells)

                            tables_data.append({
                                "index": len(tables_data),
                                "slide": slide_num,
                                "rows": rows_data,
                                "rowCount": len(rows_data),
                                "columnCount": len(rows_data[0]) if rows_data else 0,
                                "caption": "",
                            })

                            body_html_parts.append("  <table>")
                            for r_idx, row_cells in enumerate(rows_data):
                                body_html_parts.append("    <tr>")
                                cell_tag = "th" if r_idx == 0 else "td"
                                for cell in row_cells:
                                    body_html_parts.append(
                                        f"      <{cell_tag}>{_escape(cell)}</{cell_tag}>"
                                    )
                                body_html_parts.append("    </tr>")
                            body_html_parts.append("  </table>")
                        except Exception:
                            logger.debug("Failed to extract table from slide %d", slide_num)

                except Exception:
                    logger.debug("Failed to process shape on slide %d", slide_num)

            # Record slide heading
            if slide_title:
                headings.append({
                    "tag": "h2",
                    "level": 2,
                    "text": slide_title,
                    "slide": slide_num,
                })
                body_html_parts.insert(
                    len(body_html_parts) - (len(body_html_parts) - body_html_parts.index(
                        f'<section aria-label="Slide {slide_num}">'
                    ) - 1),
                    f"  <h2>{_escape(slide_title)}</h2>",
                )

            body_html_parts.append("</section>")

        pseudo_html = _wrap_html(title, body_html_parts, links, lang=doc_language)
        capture_data.html = pseudo_html
        capture_data.headings = headings
        capture_data.links = links
        capture_data.images = images
        capture_data.tables = tables_data

    except Exception:
        logger.exception("PPTX capture failed for %s", file_path)

    return capture_data


# ═════════════════════════════════════════════════════════════════════════════
#  Shared helpers
# ═════════════════════════════════════════════════════════════════════════════

def _escape(text: str) -> str:
    """Minimal HTML entity escaping."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _wrap_html(
    title: str,
    body_parts: list[str],
    links: list[dict] | None = None,
    lang: str = "en",
) -> str:
    """Wrap body content in a complete pseudo-HTML document."""
    parts = [
        "<!DOCTYPE html>",
        f'<html lang="{_escape(lang)}">',
        "<head>",
        f"  <title>{_escape(title)}</title>",
        "</head>",
        "<body>",
    ]
    parts.extend(body_parts)

    # Append link list if provided
    if links:
        parts.append('  <nav aria-label="Links">')
        for link in links:
            href = _escape(link.get("href", ""))
            text = _escape(link.get("text", "") or href)
            parts.append(f'    <a href="{href}">{text}</a>')
        parts.append("  </nav>")

    parts.extend(["</body>", "</html>"])
    return "\n".join(parts)
